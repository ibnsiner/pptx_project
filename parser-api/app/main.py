import base64
import hashlib
import io
import os
import posixpath
import re
import uuid
import xml.etree.ElementTree as ET
import zipfile
from typing import Any

from fastapi import FastAPI, File, HTTPException, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.oxml.ns import qn
from pptx.shapes.group import GroupShape

from app.slide_raster import render_slide_rasters_jpeg
from app.slide_raster_ppt import render_slide_rasters_ppt

app = FastAPI(title="PPTX Parser API")

PARSER_API_BUILD = "2026-04-01-table-element"

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


def _pct(value: int, total: int) -> str:
    if total <= 0:
        return "0%"
    return f"{(value / total) * 100:.4f}%"


def _shape_type_safe(shape: Any) -> Any:
    try:
        return shape.shape_type
    except (NotImplementedError, ValueError, AttributeError):
        return None


def _get_grp_transform(grp_shape: Any) -> tuple[int, int, float, float, int, int]:
    """
    그룹 도형의 좌표 변환 파라미터를 반환한다.
    Returns (off_x, off_y, scale_x, scale_y, ch_off_x, ch_off_y)

    PPTX 그룹 내부 좌표(chOff/chExt 기준)를 슬라이드 절대 EMU로 변환하려면:
      slide_x = off_x + (child_x - ch_off_x) * scale_x
      slide_y = off_y + (child_y - ch_off_y) * scale_y
    """
    try:
        elm = grp_shape.element
        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
        NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

        grp_sp_pr = elm.find(f"{{{NS_P}}}grpSpPr")
        if grp_sp_pr is None:
            # drawingml 네임스페이스 fallback
            grp_sp_pr = elm.find(f"{{{NS_A}}}grpSpPr")
        if grp_sp_pr is None:
            for c in elm:
                if c.tag.endswith("}grpSpPr"):
                    grp_sp_pr = c
                    break

        if grp_sp_pr is None:
            raise ValueError("no grpSpPr")

        xfrm = None
        for c in grp_sp_pr:
            if c.tag.endswith("}xfrm"):
                xfrm = c
                break
        if xfrm is None:
            raise ValueError("no xfrm")

        off = xfrm.find(f"{{{NS_A}}}off")
        ext = xfrm.find(f"{{{NS_A}}}ext")
        ch_off = xfrm.find(f"{{{NS_A}}}chOff")
        ch_ext = xfrm.find(f"{{{NS_A}}}chExt")

        off_x = int(off.get("x", 0)) if off is not None else 0
        off_y = int(off.get("y", 0)) if off is not None else 0
        ext_cx = int(ext.get("cx", 1)) if ext is not None else 1
        ext_cy = int(ext.get("cy", 1)) if ext is not None else 1
        ch_off_x = int(ch_off.get("x", 0)) if ch_off is not None else 0
        ch_off_y = int(ch_off.get("y", 0)) if ch_off is not None else 0
        ch_ext_cx = int(ch_ext.get("cx", ext_cx)) if ch_ext is not None else ext_cx
        ch_ext_cy = int(ch_ext.get("cy", ext_cy)) if ch_ext is not None else ext_cy

        scale_x = ext_cx / ch_ext_cx if ch_ext_cx != 0 else 1.0
        scale_y = ext_cy / ch_ext_cy if ch_ext_cy != 0 else 1.0
        return off_x, off_y, scale_x, scale_y, ch_off_x, ch_off_y
    except Exception:
        # 실패 시 단순 오프셋 방식 폴백
        gl = int(getattr(grp_shape, "left", 0) or 0)
        gt = int(getattr(grp_shape, "top", 0) or 0)
        return gl, gt, 1.0, 1.0, 0, 0


def _iter_shapes_placed(
    shapes: Any,
    off_x: int = 0,
    off_y: int = 0,
    scale_x: float = 1.0,
    scale_y: float = 1.0,
    ch_off_x: int = 0,
    ch_off_y: int = 0,
) -> Any:
    """
    Yield (shape, abs_slide_left_emu, abs_slide_top_emu).

    그룹 도형은 chOff/chExt 기반 좌표 변환을 적용해
    자식 도형의 위치를 슬라이드 절대 EMU로 변환한다.
    """
    for shape in shapes:
        st = _shape_type_safe(shape)
        if st == MSO_SHAPE_TYPE.GROUP or isinstance(shape, GroupShape):
            g_off_x, g_off_y, g_sx, g_sy, g_ch_x, g_ch_y = _get_grp_transform(shape)
            # 이 그룹의 슬라이드 절대 위치 (부모 변환 적용)
            child_raw_x = int(getattr(shape, "left", 0) or 0)
            child_raw_y = int(getattr(shape, "top", 0) or 0)
            slide_g_x = off_x + (child_raw_x - ch_off_x) * scale_x
            slide_g_y = off_y + (child_raw_y - ch_off_y) * scale_y

            # 하위 그룹 변환: 이 그룹의 ext/chExt 비율로 스케일 계산
            try:
                g_ext_cx = int(getattr(shape, "width", 1) or 1)
                g_ext_cy = int(getattr(shape, "height", 1) or 1)
            except Exception:
                g_ext_cx, g_ext_cy = 1, 1
            g_new_sx = scale_x * (g_ext_cx / g_ch_y if g_ch_y else 1.0)
            # grpSpPr에서 직접 chExt 읽기
            try:
                _, _, grp_sx2, grp_sy2, grp_chx2, grp_chy2 = _get_grp_transform(shape)
                new_sx = grp_sx2
                new_sy = grp_sy2
                new_chx = grp_chx2
                new_chy = grp_chy2
                new_ox = g_off_x
                new_oy = g_off_y
            except Exception:
                new_ox, new_oy = int(slide_g_x), int(slide_g_y)
                new_sx, new_sy = scale_x, scale_y
                new_chx, new_chy = ch_off_x, ch_off_y

            inner = getattr(shape, "shapes", None)
            if inner is not None:
                yield from _iter_shapes_placed(inner, new_ox, new_oy, new_sx, new_sy, new_chx, new_chy)
        else:
            # 현재 변환 적용해 슬라이드 절대 EMU 계산
            try:
                raw_x = int(shape.left or 0)
                raw_y = int(shape.top or 0)
            except Exception:
                raw_x, raw_y = 0, 0
            abs_x = off_x + (raw_x - ch_off_x) * scale_x
            abs_y = off_y + (raw_y - ch_off_y) * scale_y
            yield shape, int(abs_x), int(abs_y)


def _is_placeholder_shape(shape: Any) -> bool:
    try:
        return bool(getattr(shape, "is_placeholder", False))
    except Exception:
        return False


def _should_skip_master_layout_shape(shape: Any) -> bool:
    """
    마스터/레이아웃에서 대부분의 플레이스홀더는 슬라이드에 복제되므로 건너뛴다.
    그림·클립아트·슬라이드 이미지 플레이스홀더는 슬라이드 XML에 없을 수 있어 포함한다.
    """
    if not _is_placeholder_shape(shape):
        return False
    try:
        t = shape.placeholder_format.type
        if t in (
            PP_PLACEHOLDER.PICTURE,
            PP_PLACEHOLDER.BITMAP,
            PP_PLACEHOLDER.SLIDE_IMAGE,
            PP_PLACEHOLDER.TABLE,
        ):
            return False
    except Exception:
        pass
    return True


def _iter_shapes_paint_order(slide: Any) -> Any:
    """
    Bottom-to-top painting order (matches typical PPT stacking).
    Shapes that exist only on slide master / layout (not copied to slide XML)
    are invisible to slide.shapes; include non-placeholder shapes from those parts.
    Placeholders are skipped on master/layout because the slide carries filled copies.
    """
    try:
        layout = slide.slide_layout
        master = layout.slide_master
    except Exception:
        yield from _iter_shapes_placed(slide.shapes)
        return

    for tree in (master.shapes, layout.shapes):
        for shape, abs_x, abs_y in _iter_shapes_placed(tree):
            if not _should_skip_master_layout_shape(shape):
                yield shape, abs_x, abs_y

    yield from _iter_shapes_placed(slide.shapes)


def _upload_image_if_configured(
    content: bytes,
    content_type: str,
    slide_num: int,
    index: int,
) -> str | None:
    import os

    url = os.getenv("SUPABASE_URL")
    key = os.getenv("SUPABASE_SERVICE_ROLE_KEY")
    bucket = os.getenv("STORAGE_BUCKET", "lecture-assets")
    if not url or not key:
        return None
    try:
        from supabase import create_client

        client = create_client(url, key)
        ext = "png"
        if "jpeg" in content_type or "jpg" in content_type:
            ext = "jpg"
        path = f"draft/{uuid.uuid4().hex}/s{slide_num}_{index}.{ext}"
        client.storage.from_(bucket).upload(
            path,
            content,
            file_options={"content-type": content_type or "application/octet-stream"},
        )
        pub = client.storage.from_(bucket).get_public_url(path)
        return pub
    except Exception:
        return None


def _blip_rel_ids(shape_elm: Any) -> list[str]:
    """Collect r:embed and r:link from a:blip (linked or embedded pictures)."""
    seen: list[str] = []
    found: set[str] = set()
    q_embed = qn("r:embed")
    q_link = qn("r:link")
    for node in shape_elm.iter():
        if node.tag != qn("a:blip"):
            continue
        for attr in (q_embed, q_link):
            rid = node.get(attr)
            if rid and rid not in found:
                found.add(rid)
                seen.append(rid)
    return seen


def _load_related_image_any(part: Any, slide_part: Any, r_id: str) -> tuple[bytes | None, str | None]:
    for p in (part, slide_part):
        if p is None:
            continue
        try:
            related = p.related_part(r_id)
            blob = related.blob
            ct = (getattr(related, "content_type", None) or "").strip()
            return blob, ct or None
        except Exception:
            continue
    return None, None


def _looks_like_raster(blob: bytes | None, content_type: str | None) -> bool:
    if not blob or len(blob) < 4:
        return False
    if content_type and content_type.lower().startswith("image/"):
        return True
    if blob.startswith(b"\xff\xd8\xff"):
        return True
    if blob.startswith(b"\x89PNG\r\n\x1a\n"):
        return True
    if blob[:4] in (b"GIF8", b"RIFF"):
        return True
    if blob.startswith(b"BM"):
        return True
    if len(blob) >= 12 and blob[:4] == b"RIFF" and blob[8:12] == b"WEBP":
        return True
    # EMF: 첫 레코드 EMR_HEADER(type=1, size>=88), placeable WMF
    if len(blob) >= 8:
        rtype = int.from_bytes(blob[0:4], "little")
        rsize = int.from_bytes(blob[4:8], "little")
        if rtype == 1 and 88 <= rsize <= len(blob):
            return True
    if blob.startswith(b"\xd7\xcd\xc6\x9a"):
        return True
    return False


def _table_text_from_tbl_xml(shape_elm: Any) -> str:
    """
    a:tbl에서 텍스트 추출.
    a:tr / a:tc 구조가 있으면 행·열 순서 유지, 없으면 a:t 순서 폴백.
    """
    q_tbl = qn("a:tbl")
    q_tr = qn("a:tr")
    q_tc = qn("a:tc")
    q_t = qn("a:t")
    tbl_blocks: list[str] = []

    for tbl in shape_elm.iter():
        if tbl.tag != q_tbl:
            continue
        trs = tbl.findall(q_tr)
        if trs:
            row_strs: list[str] = []
            for tr in trs:
                cells: list[str] = []
                for tc in tr.findall(q_tc):
                    chunks: list[str] = []
                    for node in tc.iter():
                        if node.tag == q_t and node.text:
                            chunks.append(node.text)
                    cells.append("".join(chunks).strip())
                row_strs.append(" | ".join(cells))
            tbl_blocks.append("\n".join(row_strs))
        else:
            flat: list[str] = []
            for node in tbl.iter():
                if node.tag == q_t and node.text and node.text.strip():
                    flat.append(node.text.strip())
            if flat:
                tbl_blocks.append("\n".join(flat))
    return "\n\n".join(tbl_blocks)


def _text_outside_tbl(shape_elm: Any) -> str:
    skip: set[Any] = set()
    for tbl in shape_elm.iter():
        if tbl.tag == qn("a:tbl"):
            skip.update(tbl.iter())
    parts: list[str] = []
    for node in shape_elm.iter():
        if node in skip:
            continue
        if node.tag == qn("a:t") and node.text:
            parts.append(node.text)
    return "\n".join(parts)


def _inline_image_max_bytes() -> int:
    raw = os.getenv("PPTX_IMAGE_INLINE_MAX_BYTES", "").strip()
    if raw.isdigit():
        return max(1_000_000, int(raw))
    return 12_000_000


def _blob_to_data_uri(blob: bytes, content_type: str, max_bytes: int | None = None) -> str | None:
    lim = max_bytes if max_bytes is not None else _inline_image_max_bytes()
    if not blob or len(blob) > lim:
        return None
    mime = (content_type or "image/png").split(";")[0].strip() or "image/png"
    if not mime.startswith("image/"):
        mime = "image/png"
    b64 = base64.b64encode(blob).decode("ascii")
    return f"data:{mime};base64,{b64}"


def _table_text(shape: Any) -> str:
    rows_out: list[str] = []
    try:
        if not getattr(shape, "has_table", False):
            return ""
        for row in shape.table.rows:  # type: ignore[union-attr]
            cells: list[str] = []
            for cell in row.cells:
                try:
                    raw = (cell.text_frame.text or "").strip()
                    if not raw:
                        raw = (getattr(cell, "text", "") or "").strip()
                    cells.append(raw)
                except Exception:
                    try:
                        cells.append((getattr(cell, "text", "") or "").strip())
                    except Exception:
                        cells.append("")
            rows_out.append(" | ".join(cells))
    except Exception:
        try:
            el = getattr(shape, "element", None)
            if el is not None:
                return _table_text_from_tbl_xml(el)
        except Exception:
            pass
        return ""
    return "\n".join(rows_out)


def _shape_text(shape: Any) -> str:
    if getattr(shape, "has_table", False):
        t = _table_text(shape)
        if t.strip():
            return t

    t_tbl = _table_text_from_tbl_xml(shape.element)
    if t_tbl.strip():
        return t_tbl

    if getattr(shape, "has_text_frame", False):
        try:
            t = shape.text or ""
            if t.strip():
                return t
        except Exception:
            pass

    t2 = _text_outside_tbl(shape.element)
    if t2.strip():
        return t2

    return ""


def _rgb_to_hex(rgb: Any) -> str | None:
    """pptx RGBColor → '#RRGGBB' 문자열. 실패 시 None."""
    try:
        if rgb is None:
            return None
        r = int(getattr(rgb, "r", 0))
        g = int(getattr(rgb, "g", 0))
        b = int(getattr(rgb, "b", 0))
        return f"#{r:02x}{g:02x}{b:02x}"
    except Exception:
        return None


def _xml_attr_first(elm: Any, *attr_names: str) -> str | None:
    """ElementTree 노드에서 여러 속성명 중 첫 번째로 존재하는 값을 반환."""
    if elm is None:
        return None
    for a in attr_names:
        v = elm.get(a)
        if v is not None:
            return v
    return None


def _extract_paragraph_styles(shape: Any, slide_w: int) -> list[dict[str, Any]]:
    """
    TextFrame의 각 단락(paragraph)별 스타일을 리스트로 반환한다.
    단락 사이 빈 줄 포함, 각 단락의 첫 런 스타일을 대표로 사용한다.
    빈 단락은 빈 딕셔너리로 표시.
    """
    result: list[dict[str, Any]] = []
    try:
        tf = getattr(shape, "text_frame", None)
        if tf is None:
            return result
        for para in tf.paragraphs:
            # 단락 텍스트
            try:
                para_text = para.text or ""
            except Exception:
                para_text = ""
            if not para_text.strip():
                result.append({})
                continue
            # 단락 스타일: _scan_txbody 대신 직접 rPr 탐색
            para_style: dict[str, Any] = {}
            slide_w_pt = float(slide_w) / 12700.0
            NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
            def _tag(local: str) -> str:
                return f"{{{NS_A}}}{local}"
            try:
                p_elm = para._p
                # textAlign
                pPr = p_elm.find(_tag("pPr"))
                if pPr is not None:
                    algn = pPr.get("algn")
                    am = {"l": "left", "ctr": "center", "r": "right", "just": "justify"}
                    if algn in am:
                        para_style["textAlign"] = am[algn]
                # 런 스타일
                for r_elm in p_elm.findall(_tag("r")):
                    rpr = r_elm.find(_tag("rPr"))
                    if rpr is None:
                        continue
                    b = rpr.get("b")
                    if b is not None and "bold" not in para_style:
                        para_style["bold"] = b not in ("0", "false")
                    i = rpr.get("i")
                    if i is not None and "italic" not in para_style:
                        para_style["italic"] = i not in ("0", "false")
                    sz = rpr.get("sz")
                    if sz and "fontSize" not in para_style and slide_w_pt > 0:
                        try:
                            cqw = min(max((int(sz) * 0.01 / slide_w_pt) * 100.0, 0.4), 14.0)
                            para_style["fontSize"] = f"{cqw:.4f}cqw"
                        except Exception:
                            pass
                    solid = rpr.find(_tag("solidFill"))
                    if solid is not None and "color" not in para_style:
                        srgb = solid.find(_tag("srgbClr"))
                        if srgb is not None:
                            val = srgb.get("val", "")
                            if len(val) == 6:
                                para_style["color"] = f"#{val.lower()}"
                    for ftag in (_tag("latin"), _tag("ea")):
                        node = rpr.find(ftag)
                        if node is not None and "fontFamily" not in para_style:
                            tf2 = node.get("typeface", "")
                            if tf2 and not tf2.startswith("+"):
                                para_style["fontFamily"] = tf2
                    if len(para_style) >= 5:
                        break
            except Exception:
                pass
            result.append(para_style)
    except Exception:
        pass
    return result


def _extract_text_style(shape: Any, slide_w: int) -> dict[str, Any]:
    """
    텍스트 스타일을 슬라이드 → 레이아웃 → 마스터 상속 체인을 따라 추출한다.

    Placeholder 도형은 슬라이드 XML에 <a:rPr>가 없고 Layout/Master에 스타일이 정의되므로
    세 단계를 순서대로 탐색해 빈 필드를 채운다.

    반환 키: color, bold, italic, underline, fontSize, fontFamily, textAlign
    """
    style: dict[str, Any] = {}
    try:
        slide_w_pt = float(slide_w) / 12700.0
        if slide_w_pt <= 0:
            return style

        NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

        def _tag(local: str) -> str:
            return f"{{{NS_A}}}{local}"

        def _parse_rpr(rpr: Any) -> dict[str, Any]:
            """<a:rPr> 또는 <a:defRPr> XML에서 스타일 딕셔너리 생성."""
            out: dict[str, Any] = {}
            if rpr is None:
                return out
            b = rpr.get("b")
            if b is not None:
                out["bold"] = b not in ("0", "false")
            i = rpr.get("i")
            if i is not None:
                out["italic"] = i not in ("0", "false")
            u = rpr.get("u")
            if u is not None and u not in ("none", "0", "false"):
                out["underline"] = True
            sz = rpr.get("sz")
            if sz:
                try:
                    pt_val = int(sz) * 0.01
                    cqw = min(max((pt_val / slide_w_pt) * 100.0, 0.4), 14.0)
                    out["fontSize"] = f"{cqw:.4f}cqw"
                except (ValueError, ZeroDivisionError):
                    pass
            for tag in (_tag("latin"), _tag("ea")):
                node = rpr.find(tag)
                if node is not None:
                    tf = node.get("typeface", "")
                    if tf and not tf.startswith("+"):
                        out["fontFamily"] = tf
                        break
            solid = rpr.find(_tag("solidFill"))
            if solid is not None:
                srgb = solid.find(_tag("srgbClr"))
                if srgb is not None:
                    val = srgb.get("val", "")
                    if len(val) == 6:
                        out["color"] = f"#{val.lower()}"
            return out

        def _merge(base: dict, new: dict) -> None:
            for k, v in new.items():
                if k not in base:
                    base[k] = v

        def _scan_txbody(txBody: Any) -> None:
            """txBody 전체를 순회해 style에 채워 넣는다."""
            if txBody is None:
                return
            for p_elm in txBody.findall(_tag("p")):
                # textAlign
                if "textAlign" not in style:
                    pPr = p_elm.find(_tag("pPr"))
                    if pPr is not None:
                        algn = pPr.get("algn")
                        am = {"l": "left", "ctr": "center", "r": "right", "just": "justify"}
                        if algn in am:
                            style["textAlign"] = am[algn]

                # <a:r><a:rPr>
                for r_elm in p_elm.findall(_tag("r")):
                    rpr = r_elm.find(_tag("rPr"))
                    if rpr is not None:
                        _merge(style, _parse_rpr(rpr))

                # <a:br><a:rPr>
                for br_elm in p_elm.findall(_tag("br")):
                    rpr = br_elm.find(_tag("rPr"))
                    if rpr is not None:
                        _merge(style, _parse_rpr(rpr))

                # <a:pPr><a:defRPr>
                pPr = p_elm.find(_tag("pPr"))
                if pPr is not None:
                    defRPr = pPr.find(_tag("defRPr"))
                    if defRPr is not None:
                        _merge(style, _parse_rpr(defRPr))

                if len(style) >= 6:
                    return

            # lstStyle / defPPr / defRPr 폴백
            lst = txBody.find(_tag("lstStyle"))
            if lst is not None:
                for defPPr in (lst.find(_tag("defPPr")), lst.find(_tag("lvl1pPr"))):
                    if defPPr is not None:
                        defRPr = defPPr.find(_tag("defRPr"))
                        if defRPr is not None:
                            _merge(style, _parse_rpr(defRPr))
                            # textAlign from defPPr
                            if "textAlign" not in style:
                                algn = defPPr.get("algn")
                                am = {"l": "left", "ctr": "center", "r": "right", "just": "justify"}
                                if algn in am:
                                    style["textAlign"] = am[algn]

        # 1-A. 표(Table) 도형: 첫 번째 셀의 text_frame에서 대표 스타일 추출
        if getattr(shape, "has_table", False):
            try:
                for row in shape.table.rows:
                    for cell in row.cells:
                        cell_tf = getattr(cell, "text_frame", None)
                        if cell_tf is None:
                            continue
                        cell_txBody = getattr(cell_tf, "_txBody", None) or getattr(cell_tf, "_element", None)
                        if cell_txBody is not None:
                            _scan_txbody(cell_txBody)
                        if len(style) >= 3:
                            break
                    if len(style) >= 3:
                        break
            except Exception:
                pass
            return style

        # 1-B. 일반 도형 / 플레이스홀더
        tf = getattr(shape, "text_frame", None)
        if tf is None:
            return style
        txBody = getattr(tf, "_txBody", None) or getattr(tf, "_element", None)
        _scan_txbody(txBody)

        # 2. Layout Placeholder → 3. Master Placeholder (슬라이드에서 못 찾은 속성 보완)
        if len(style) < 4 and getattr(shape, "is_placeholder", False):
            ph_idx = None
            try:
                ph_idx = shape.placeholder_format.idx
            except Exception:
                pass

            sources: list[Any] = []
            try:
                layout = shape.part.slide.slide_layout
                sources.append(layout)
                sources.append(layout.slide_master)
            except Exception:
                pass

            for src in sources:
                if len(style) >= 6:
                    break
                try:
                    phs = list(src.placeholders)
                except Exception:
                    phs = []
                for ph in phs:
                    if ph_idx is not None and getattr(getattr(ph, "placeholder_format", None), "idx", None) != ph_idx:
                        continue
                    try:
                        ph_tf = ph.text_frame
                        ph_txBody = getattr(ph_tf, "_txBody", None) or getattr(ph_tf, "_element", None)
                        _scan_txbody(ph_txBody)
                    except Exception:
                        pass
                    break  # 첫 번째 매칭 ph만

    except Exception:
        pass
    return style


def _pic_geom_emu(pic_elm: Any) -> tuple[int, int, int, int] | None:
    """p:pic 아래 p:spPr/a:xfrm 에서 위치·크기(EMU)."""
    sp_pr = pic_elm.find(qn("p:spPr"))
    if sp_pr is None:
        return None
    xfrm = sp_pr.find(qn("a:xfrm"))
    if xfrm is None:
        return None
    off = xfrm.find(qn("a:off"))
    ext = xfrm.find(qn("a:ext"))
    if off is None or ext is None:
        return None
    try:
        x = int(off.get("x", 0))
        y = int(off.get("y", 0))
        cx = int(ext.get("cx", 0))
        cy = int(ext.get("cy", 0))
        if cx <= 0 or cy <= 0:
            return None
        return x, y, cx, cy
    except (TypeError, ValueError):
        return None


def _pic_blip_rel_ids(pic_elm: Any) -> list[str]:
    """p:blipFill 안의 a:blip r:embed / r:link."""
    out: list[str] = []
    found: set[str] = set()
    bf = pic_elm.find(qn("p:blipFill"))
    if bf is None:
        return out
    q_embed = qn("r:embed")
    q_link = qn("r:link")
    for blip in bf.iter():
        if blip.tag != qn("a:blip"):
            continue
        for attr in (q_embed, q_link):
            rid = blip.get(attr)
            if rid and rid not in found:
                found.add(rid)
                out.append(rid)
    return out


def _image_blob_fingerprint(blob: bytes) -> str:
    return hashlib.sha256(blob).hexdigest()


def _append_raster_dims(
    elements: list[dict[str, Any]],
    blob: bytes,
    ct: str,
    abs_l: int,
    abs_t: int,
    w_emu: int,
    h_emu: int,
    slide_w: int,
    slide_h: int,
    slide_num: int,
    img_idx: int,
    image_fps: set[str] | None = None,
) -> int:
    """같은 슬라이드에서 동일 바이너리가 여러 경로로 들어올 때 한 번만 표시."""
    if image_fps is not None and blob:
        fp = _image_blob_fingerprint(blob)
        if fp in image_fps:
            return img_idx
        image_fps.add(fp)
    left_pct = _pct(abs_l, slide_w)
    top_pct = _pct(abs_t, slide_h)
    width_pct = _pct(w_emu, slide_w)
    height_pct = _pct(h_emu, slide_h)
    url = _upload_image_if_configured(blob, ct, slide_num, img_idx)
    if not url:
        url = _blob_to_data_uri(blob, ct or "image/png")
    if not url:
        elements.append(
            {
                "type": "image",
                "src": "",
                "skipReason": "too_large_for_inline",
                "byteLength": len(blob),
                "contentType": ct or "",
                "style": {
                    "left": left_pct,
                    "top": top_pct,
                    "width": width_pct,
                    "height": height_pct,
                },
            }
        )
        return img_idx + 1
    elements.append(
        {
            "type": "image",
            "src": url,
            "style": {
                "left": left_pct,
                "top": top_pct,
                "width": width_pct,
                "height": height_pct,
            },
        }
    )
    return img_idx + 1


def _append_raster_element(
    elements: list[dict[str, Any]],
    blob: bytes,
    ct: str,
    abs_l: int,
    abs_t: int,
    shape: Any,
    slide_w: int,
    slide_h: int,
    slide_num: int,
    img_idx: int,
    image_fps: set[str] | None = None,
) -> int:
    return _append_raster_dims(
        elements,
        blob,
        ct,
        abs_l,
        abs_t,
        int(shape.width),
        int(shape.height),
        slide_w,
        slide_h,
        slide_num,
        img_idx,
        image_fps,
    )


def _xfrm_off_ext_emu(xfrm: Any) -> tuple[int, int, int, int] | None:
    if xfrm is None:
        return None
    off = xfrm.find(qn("a:off"))
    ext = xfrm.find(qn("a:ext"))
    if off is None or ext is None:
        return None
    try:
        x = int(off.get("x", 0))
        y = int(off.get("y", 0))
        cx = int(ext.get("cx", 0))
        cy = int(ext.get("cy", 0))
        if cx <= 0 or cy <= 0:
            return None
        return x, y, cx, cy
    except (TypeError, ValueError):
        return None


def _graphic_frame_geom_emu(gf_elm: Any) -> tuple[int, int, int, int] | None:
    """p:graphicFrame 직속 p:xfrm."""
    xfrm = gf_elm.find(qn("p:xfrm"))
    return _xfrm_off_ext_emu(xfrm)


def _normalized_text_block(s: str) -> str:
    return "\n".join(
        ln.strip() for ln in s.strip().splitlines() if ln.strip()
    )


def _collect_existing_text_keys(elements: list[dict[str, Any]]) -> set[str]:
    keys: set[str] = set()
    for el in elements:
        if el.get("type") != "text":
            continue
        c = (el.get("content") or "").strip()
        if c:
            keys.add(_normalized_text_block(c))
            keys.add(c)
    return keys


def _get_theme_elm(shape: Any) -> Any:
    """
    shape가 슬라이드·레이아웃·마스터 어느 파트에 속해 있어도
    해당 슬라이드 마스터의 theme XML 엘리먼트를 반환한다.
    실패 시 None.
    """
    THEME_REL = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme"
    candidates = []
    try:
        part = shape.part
        candidates.append(part)
        # slide → layout → master
        for getter in ("slide", "slide_layout"):
            try:
                part = getattr(part, getter, None) or getattr(part, getter.replace("_", "") if "_" in getter else getter, None)
                if part is None:
                    break
                candidates.append(part)
            except Exception:
                break
    except Exception:
        pass

    # 여러 경로로 master까지 올라가며 theme rel 시도
    try:
        # 직접 경로: shape.part.slide.slide_layout.slide_master.part
        master_part = shape.part.slide.slide_layout.slide_master.part
        return master_part.part_related_by(THEME_REL)._element
    except Exception:
        pass
    try:
        # layout 경로
        master_part = shape.part.slide_layout.slide_master.part
        return master_part.part_related_by(THEME_REL)._element
    except Exception:
        pass
    try:
        # shape.part 자체가 master part인 경우
        return shape.part.part_related_by(THEME_REL)._element
    except Exception:
        pass
    try:
        # shape.part.slide_master (layout shape)
        master_part = shape.part.slide_master.part
        return master_part.part_related_by(THEME_REL)._element
    except Exception:
        pass
    return None


def _resolve_scheme_color(scheme_name: str, theme_elm: Any, NS_A: str) -> str | None:
    """
    테마 XML에서 schemeClr 이름(예: accent1, dk1 등)을 실제 RGB hex로 변환한다.
    theme_elm: pptx의 theme1.xml 루트 엘리먼트 (없으면 None)
    """
    if theme_elm is None:
        return None

    def _tag(local: str) -> str:
        return f"{{{NS_A}}}{local}"

    NS_THM = "http://schemas.openxmlformats.org/drawingml/2006/main"

    # themeElements/clrScheme 탐색
    theme_elems = None
    for child in theme_elm:
        if child.tag.endswith("}themeElements"):
            theme_elems = child
            break
    if theme_elems is None:
        return None
    clr_scheme = None
    for child in theme_elems:
        if child.tag.endswith("}clrScheme"):
            clr_scheme = child
            break
    if clr_scheme is None:
        return None

    # scheme_name → 엘리먼트 이름 매핑
    name_map = {
        "dk1": "dk1", "lt1": "lt1", "dk2": "dk2", "lt2": "lt2",
        "accent1": "accent1", "accent2": "accent2", "accent3": "accent3",
        "accent4": "accent4", "accent5": "accent5", "accent6": "accent6",
        "hlink": "hlink", "folHlink": "folHlink",
    }
    target = name_map.get(scheme_name)
    if not target:
        return None

    for clr_elm in clr_scheme:
        local = clr_elm.tag.split("}")[-1] if "}" in clr_elm.tag else clr_elm.tag
        if local != target:
            continue
        for color_node in clr_elm:
            local2 = color_node.tag.split("}")[-1]
            if local2 == "srgbClr":
                val = color_node.get("val", "")
                if len(val) == 6:
                    return f"#{val.lower()}"
            elif local2 == "sysClr":
                last_clr = color_node.get("lastClr", "")
                if len(last_clr) == 6:
                    return f"#{last_clr.lower()}"
    return None


def _extract_line_color_from_sppr(
    sp_pr: Any, NS_A: str, theme_elm: Any = None
) -> str | None:
    """<a:spPr> 또는 <a:cxnSpPr>에서 선 색상을 추출한다."""
    def _tag(local: str) -> str:
        return f"{{{NS_A}}}{local}"
    ln = sp_pr.find(_tag("ln"))
    if ln is None:
        return None
    solid = ln.find(_tag("solidFill"))
    if solid is None:
        return None
    srgb = solid.find(_tag("srgbClr"))
    if srgb is not None:
        val = srgb.get("val", "")
        if len(val) == 6:
            return f"#{val.lower()}"
    # schemeClr: 테마 XML로 실제 색상 해석 시도
    schm = solid.find(_tag("schemeClr"))
    if schm is not None:
        name = schm.get("val", "")
        resolved = _resolve_scheme_color(name, theme_elm, NS_A)
        if resolved:
            return resolved
        return f"scheme:{name}"
    return None


def _try_extract_connector_shape(
    shape: Any,
    abs_l: int,
    abs_t: int,
    slide_w: int,
    slide_h: int,
    elements: list[dict[str, Any]],
    global_theme_elm: Any = None,
) -> None:
    """
    <p:cxnSp> 연결선·화살표·장식선을 type:"shape"으로 추출한다.
    HTML에서는 얇은 색상 막대로 렌더링한다.
    """
    try:
        sp_elm = getattr(shape, "element", None)
        if sp_elm is None:
            return

        NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

        def _tag_p(local: str) -> str:
            return f"{{{NS_P}}}{local}"

        # p:cxnSp 인지 확인
        is_cxnSp = sp_elm.tag.endswith("}cxnSp") or sp_elm.tag == _tag_p("cxnSp")
        if not is_cxnSp:
            return

        # 크기·위치
        try:
            w = shape.width
            h = shape.height
        except Exception:
            return
        if w <= 0 and h <= 0:
            return

        # cxnSpPr에서 선 색상 추출 (테마 XML 포함)
        sp_pr = None
        for child in sp_elm:
            if child.tag.endswith("}cxnSpPr") or child.tag.endswith("}spPr"):
                sp_pr = child
                break
        if sp_pr is None:
            return

        # 테마 XML 획득: ZIP 직접 추출 > part 체인
        theme_elm = global_theme_elm or _get_theme_elm(shape)

        line_color = _extract_line_color_from_sppr(sp_pr, NS_A, theme_elm)
        if not line_color or line_color.startswith("scheme:"):
            line_color = "#888888" if line_color else None
        if not line_color:
            return

        # 가로/세로 방향에 따라 두께 결정
        thickness = max(h, 1) if w > h else max(w, 1)
        is_horizontal = w >= h

        elements.append(
            {
                "type": "shape",
                "style": {
                    "left": _pct(abs_l, slide_w),
                    "top": _pct(abs_t, slide_h),
                    "width": _pct(w, slide_w) if not is_horizontal else _pct(w, slide_w),
                    "height": _pct(h, slide_h),
                    "fillColor": line_color,
                    "fillOpacity": "1",
                    "borderRadius": "0%",
                    "isLine": True,
                },
            }
        )
    except Exception:
        pass


def _try_extract_filled_shape(
    shape: Any,
    abs_l: int,
    abs_t: int,
    slide_w: int,
    slide_h: int,
    elements: list[dict[str, Any]],
    global_theme_elm: Any = None,
) -> None:
    """
    텍스트 없는 채워진 도형(원, 사각형 등)을 type:"shape" 요소로 추출한다.
    HTML 렌더링에서 CSS background-color + border-radius로 표현한다.

    추출 조건: solidFill(단색) 또는 gradFill(그라데이션) 이 있는 p:sp 도형.
    마스터/레이아웃 기본 도형(흰 배경, 투명)은 제외.
    """
    try:
        # 이미지/표 도형 제외
        if getattr(shape, "has_table", False):
            return
        sp_elm = getattr(shape, "element", None)
        if sp_elm is None:
            return

        NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
        NS_P = "http://schemas.openxmlformats.org/presentationml/2006/main"

        def _tag_a(local: str) -> str:
            return f"{{{NS_A}}}{local}"

        # spPr 탐색
        sp_pr = sp_elm.find(f"{{{NS_P}}}spPr") or sp_elm.find(_tag_a("spPr"))
        if sp_pr is None:
            # p:sp 구조에서 직접 탐색
            for child in sp_elm:
                if child.tag.endswith("}spPr"):
                    sp_pr = child
                    break
        if sp_pr is None:
            return

        # 도형 종류 + flipH/flipV — prstGeom & xfrm
        border_radius = "0%"
        clip_path: str | None = None

        # xfrm에서 뒤집기 정보 추출 — 네임스페이스 무관 탐색
        xfrm = None
        for _child in sp_pr:
            local = _child.tag.split("}")[-1] if "}" in _child.tag else _child.tag
            if local == "xfrm":
                xfrm = _child
                break
        flip_h = xfrm is not None and xfrm.get("flipH") in ("1", "true")
        flip_v = xfrm is not None and xfrm.get("flipV") in ("1", "true")
        # 180도 회전(rot=10800000)도 flipV와 동일 효과
        if xfrm is not None and not flip_v:
            rot_val = xfrm.get("rot", "0")
            try:
                rot_deg = int(rot_val) / 60000.0
                if abs(rot_deg - 180.0) < 1.0:
                    flip_v = True
                elif abs(rot_deg - 90.0) < 1.0 or abs(rot_deg - 270.0) < 1.0:
                    flip_h = True
            except (ValueError, ZeroDivisionError):
                pass

        prstGeom = sp_pr.find(_tag_a("prstGeom"))
        if prstGeom is not None:
            prst = prstGeom.get("prst", "")
            if prst in ("ellipse", "circle"):
                border_radius = "50%"
            elif prst in ("roundRect", "round1Rect", "round2SameRect"):
                border_radius = "12px"
            elif prst == "triangle":
                # 기본 삼각형: 위쪽 꼭짓점 (∧)
                # flipV → 아래 꼭짓점 (∨), flipH → 좌우 반전
                if flip_v:
                    clip_path = "polygon(0% 0%, 100% 0%, 50% 100%)"
                elif flip_h:
                    clip_path = "polygon(100% 0%, 0% 50%, 100% 100%)"
                else:
                    clip_path = "polygon(50% 0%, 100% 100%, 0% 100%)"
            elif prst == "rtTriangle":
                if flip_h:
                    clip_path = "polygon(100% 0%, 0% 100%, 100% 100%)"
                else:
                    clip_path = "polygon(0% 0%, 100% 100%, 0% 100%)"
            elif prst in ("rightArrow", "leftArrow"):
                # 화살표: 단순 오각형으로 근사
                if prst == "rightArrow":
                    clip_path = "polygon(0% 25%, 70% 25%, 70% 0%, 100% 50%, 70% 100%, 70% 75%, 0% 75%)"
                else:
                    clip_path = "polygon(100% 25%, 30% 25%, 30% 0%, 0% 50%, 30% 100%, 30% 75%, 100% 75%)"
            elif prst in ("pentagon", "hexagon"):
                if prst == "pentagon":
                    clip_path = "polygon(50% 0%, 100% 38%, 82% 100%, 18% 100%, 0% 38%)"
                else:
                    clip_path = "polygon(25% 0%, 75% 0%, 100% 50%, 75% 100%, 25% 100%, 0% 50%)"
            elif prst == "diamond":
                clip_path = "polygon(50% 0%, 100% 50%, 50% 100%, 0% 50%)"
            elif prst == "parallelogram":
                clip_path = "polygon(25% 0%, 100% 0%, 75% 100%, 0% 100%)"

        # 채우기 색상 추출
        fill_color: str | None = None
        fill_opacity: float = 1.0

        # solidFill
        solid = sp_pr.find(_tag_a("solidFill"))
        if solid is not None:
            srgb = solid.find(_tag_a("srgbClr"))
            if srgb is not None:
                val = srgb.get("val", "")
                if len(val) == 6:
                    # 투명도(alpha) 확인
                    alpha_elm = srgb.find(_tag_a("alpha"))
                    if alpha_elm is not None:
                        try:
                            fill_opacity = int(alpha_elm.get("val", "100000")) / 100000.0
                        except ValueError:
                            pass
                    fill_color = f"#{val.lower()}"
            # schemeClr(테마 색상): 직접 해석 어려우므로 스킵
        else:
            # gradFill: 첫 번째 stop 색상 사용
            grad = sp_pr.find(_tag_a("gradFill"))
            if grad is not None:
                gsLst = grad.find(_tag_a("gsLst"))
                if gsLst is not None:
                    for gs in gsLst:
                        srgb = gs.find(_tag_a("srgbClr"))
                        if srgb is None:
                            # solidFill 등 내부 탐색
                            for c in gs:
                                srgb = c.find(_tag_a("srgbClr")) if c is not None else None
                                if srgb is not None:
                                    break
                        if srgb is not None:
                            val = srgb.get("val", "")
                            if len(val) == 6:
                                fill_color = f"#{val.lower()}"
                                break

        # 테마 XML 획득: ZIP 직접 추출 > part 체인
        _theme_elm = global_theme_elm or _get_theme_elm(shape)

        # schemeClr solidFill 폴백: 테마 색상으로 해석 시도
        if fill_color is None:
            solid = sp_pr.find(_tag_a("solidFill"))
            if solid is not None:
                schm = solid.find(_tag_a("schemeClr"))
                if schm is not None:
                    scheme_name = schm.get("val", "")
                    resolved = _resolve_scheme_color(scheme_name, _theme_elm, NS_A)
                    if resolved:
                        # luminance modifiers (lumMod, lumOff) 적용
                        # OOXML: lumMod/lumOff 값은 100000 = 100%
                        # lumMod: 현재 밝기에 곱함 (50000 = 50%로 어둡게)
                        # lumOff: 밝기에 더함 (50000 = +50% 밝게)
                        try:
                            import colorsys
                            children = list(schm)
                            mods: dict[str, int] = {}
                            for c in children:
                                local = c.tag.split("}")[-1] if "}" in c.tag else c.tag
                                val_str = c.get("val", "0")
                                try:
                                    mods[local] = int(val_str)
                                except ValueError:
                                    pass
                            if mods:
                                hex_val = resolved.lstrip("#")
                                r, g, b = (int(hex_val[i:i+2], 16) / 255.0 for i in (0, 2, 4))
                                h, l, s = colorsys.rgb_to_hls(r, g, b)
                                if "lumMod" in mods:
                                    l = l * mods["lumMod"] / 100000.0
                                if "lumOff" in mods:
                                    l = l + mods["lumOff"] / 100000.0
                                if "shade" in mods:
                                    l = l * mods["shade"] / 100000.0
                                if "tint" in mods:
                                    l = l + (1.0 - l) * mods["tint"] / 100000.0
                                l = max(0.0, min(1.0, l))
                                r2, g2, b2 = colorsys.hls_to_rgb(h, l, s)
                                resolved = "#{:02x}{:02x}{:02x}".format(
                                    int(r2 * 255), int(g2 * 255), int(b2 * 255)
                                )
                        except Exception:
                            pass
                        fill_color = resolved

        # noFill인지 확인: noFill이 있으면 진짜로 투명, 없으면 scheme 기반 폴백 사용
        if fill_color is None:
            has_no_fill = sp_pr.find(_tag_a("noFill")) is not None
            if not has_no_fill:
                # solidFill에 schemeClr이 있지만 해석 실패한 경우 → scheme 이름 기반 폴백
                solid2 = sp_pr.find(_tag_a("solidFill"))
                if solid2 is not None:
                    schm2 = solid2.find(_tag_a("schemeClr"))
                    if schm2 is not None:
                        sname = schm2.get("val", "")
                        # 어두운 계열 → 중간 회색 폴백
                        dark_schemes = {"dk1", "dk2", "bg2", "tx2"}
                        light_schemes = {"lt1", "lt2", "bg1", "tx1"}
                        if sname in dark_schemes:
                            fill_color = "#6b7280"
                        elif sname in light_schemes:
                            fill_color = "#d1d5db"
                        elif sname.startswith("accent"):
                            # 테마 accent 색상을 못 읽은 경우 중립 파란색
                            fill_color = "#3b82f6"

        # 채우기 없으면 stroke-only 도형 시도
        if fill_color is None:
            line_color = _extract_line_color_from_sppr(sp_pr, NS_A, _theme_elm)
            if line_color and not line_color.startswith("scheme:"):
                try:
                    w = shape.width
                    h = shape.height
                except Exception:
                    return
                if w > 0 and h > 0:
                    elements.append(
                        {
                            "type": "shape",
                            "style": {
                                "left": _pct(abs_l, slide_w),
                                "top": _pct(abs_t, slide_h),
                                "width": _pct(w, slide_w),
                                "height": _pct(h, slide_h),
                                "fillColor": "transparent",
                                "fillOpacity": "1",
                                "borderRadius": border_radius,
                                "strokeColor": line_color,
                            },
                        }
                    )
            return
        if fill_color.lower() in ("#ffffff", "#fff", "#fefefe", "#fdfdfd"):
            return
        if fill_opacity < 0.05:
            return

        try:
            width = shape.width
            height = shape.height
        except Exception:
            return
        if width <= 0 or height <= 0:
            return

        opacity_str = f"{fill_opacity:.3f}" if fill_opacity < 1.0 else "1"
        shape_style: dict[str, Any] = {
            "left": _pct(abs_l, slide_w),
            "top": _pct(abs_t, slide_h),
            "width": _pct(width, slide_w),
            "height": _pct(height, slide_h),
            "fillColor": fill_color,
            "fillOpacity": opacity_str,
            "borderRadius": border_radius,
        }
        if clip_path:
            shape_style["clipPath"] = clip_path
        # 디버깅: prst + flip 정보
        if prstGeom is not None:
            _prst = prstGeom.get("prst", "")
            if _prst:
                shape_style["_prst"] = _prst
        if xfrm is not None:
            _fh = xfrm.get("flipH")
            _fv = xfrm.get("flipV")
            _rot = xfrm.get("rot")
            if _fh:
                shape_style["_flipH"] = _fh
            if _fv:
                shape_style["_flipV"] = _fv
            if _rot:
                shape_style["_rot"] = _rot
        elements.append({"type": "shape", "style": shape_style})
    except Exception:
        pass


def _extract_text_style_from_tbl_xml(tbl_elm: Any, slide_w: int) -> dict[str, Any]:
    """
    <a:tbl> XML 엘리먼트에서 첫 번째 셀의 대표 스타일을 추출한다.
    표는 shape 객체가 없어 직접 XML 탐색이 필요하다.
    """
    style: dict[str, Any] = {}
    if tbl_elm is None or slide_w <= 0:
        return style
    try:
        slide_w_pt = float(slide_w) / 12700.0
        NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"

        def _tag(local: str) -> str:
            return f"{{{NS_A}}}{local}"

        def _parse_rpr_xml(rpr: Any) -> dict[str, Any]:
            out: dict[str, Any] = {}
            if rpr is None:
                return out
            b = rpr.get("b")
            if b is not None:
                out["bold"] = b not in ("0", "false")
            i = rpr.get("i")
            if i is not None:
                out["italic"] = i not in ("0", "false")
            sz = rpr.get("sz")
            if sz:
                try:
                    cqw = min(max((int(sz) * 0.01 / slide_w_pt) * 100.0, 0.4), 14.0)
                    out["fontSize"] = f"{cqw:.4f}cqw"
                except (ValueError, ZeroDivisionError):
                    pass
            for tag in (_tag("latin"), _tag("ea")):
                node = rpr.find(tag)
                if node is not None:
                    tf = node.get("typeface", "")
                    if tf and not tf.startswith("+"):
                        out["fontFamily"] = tf
                        break
            solid = rpr.find(_tag("solidFill"))
            if solid is not None:
                srgb = solid.find(_tag("srgbClr"))
                if srgb is not None:
                    val = srgb.get("val", "")
                    if len(val) == 6:
                        out["color"] = f"#{val.lower()}"
            return out

        # <a:tbl> → <a:tr> → <a:tc> → <a:txBody> → <a:p>
        tbl = tbl_elm.find(f"{{{NS_A}}}tbl") if tbl_elm.tag != f"{{{NS_A}}}tbl" else tbl_elm
        if tbl is None:
            return style
        for tr in tbl.findall(_tag("tr")):
            for tc in tr.findall(_tag("tc")):
                txBody = tc.find(_tag("txBody"))
                if txBody is None:
                    continue
                for p in txBody.findall(_tag("p")):
                    if "textAlign" not in style:
                        pPr = p.find(_tag("pPr"))
                        if pPr is not None:
                            algn = pPr.get("algn")
                            am = {"l": "left", "ctr": "center", "r": "right", "just": "justify"}
                            if algn in am:
                                style["textAlign"] = am[algn]
                    for r in p.findall(_tag("r")):
                        rpr = r.find(_tag("rPr"))
                        if rpr is not None:
                            for k, v in _parse_rpr_xml(rpr).items():
                                if k not in style:
                                    style[k] = v
                    if len(style) >= 4:
                        return style
                if len(style) >= 4:
                    return style
    except Exception:
        pass
    return style


def _enrich_tables_from_slide_oxml(
    slide_elm: Any,
    slide_w: int,
    slide_h: int,
    elements: list[dict[str, Any]],
    slide_text_parts: list[str],
) -> None:
    """
    slide.shapes / python-pptx에 안 올라오는 표(p:graphicFrame, mc:AlternateContent 등) 보강.
    """
    q_gf = qn("p:graphicFrame")
    seen = _collect_existing_text_keys(elements)
    for gf in slide_elm.iter():
        if gf.tag != q_gf:
            continue
        t = _table_text_from_tbl_xml(gf)
        if not t.strip():
            continue
        key = _normalized_text_block(t)
        if key in seen:
            continue
        seen.add(key)
        slide_text_parts.append(t.strip())
        geom = _graphic_frame_geom_emu(gf)
        if geom:
            abs_l, abs_t, w_emu, h_emu = geom
        else:
            abs_l, abs_t = 0, 0
            w_emu = max(slide_w // 2, 1)
            h_emu = max(slide_h // 8, 1)

        # 표 셀 스타일 추출 (XML 직접 탐색)
        tbl_style = _extract_text_style_from_tbl_xml(gf, slide_w)

        cell_style: dict[str, Any] = {
            "left": _pct(abs_l, slide_w),
            "top": _pct(abs_t, slide_h),
            "width": _pct(w_emu, slide_w),
            "height": _pct(h_emu, slide_h),
        }
        cell_style.update(tbl_style)
        elements.append(
            {
                "type": "text",
                "content": t,
                "style": cell_style,
            }
        )


def _blip_under_slide_background(blip: Any) -> bool:
    cur: Any = blip
    for _ in range(40):
        if cur is None:
            return False
        if cur.tag in (qn("p:bg"), qn("p:bgPr")):
            return True
        cur = cur.getparent()
    return False


def _geom_from_blip_ancestor(blip: Any) -> tuple[int, int, int, int] | None:
    """a:blip에서 상위 도형의 xfrm(off/ext) 추적."""
    cur: Any = blip.getparent()
    q_sp = qn("p:sp")
    q_pic = qn("p:pic")
    q_gf = qn("p:graphicFrame")
    q_grp = qn("p:grpSp")
    q_cxn = qn("p:cxnSp")
    while cur is not None:
        tag = cur.tag
        if tag == q_pic:
            g = _pic_geom_emu(cur)
            if g:
                return g
        if tag in (q_sp, q_cxn):
            sp_pr = cur.find(qn("p:spPr"))
            if sp_pr is not None:
                xf = sp_pr.find(qn("a:xfrm"))
                r = _xfrm_off_ext_emu(xf)
                if r:
                    return r
        if tag == q_gf:
            xf = cur.find(qn("p:xfrm"))
            r = _xfrm_off_ext_emu(xf)
            if r:
                return r
        if tag == q_grp:
            gsp = cur.find(qn("p:grpSpPr"))
            if gsp is not None:
                xf = gsp.find(qn("a:xfrm"))
                r = _xfrm_off_ext_emu(xf)
                if r:
                    return r
        cur = cur.getparent()
    return None


_OOXML_REL_EMBED = (
    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
)
_OOXML_REL_LINK = (
    "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}link"
)
_OOXML_REL_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"


def _xml_local_tag(el: Any) -> str:
    t = el.tag
    return t.rsplit("}", 1)[-1] if isinstance(t, str) and "}" in t else str(t)


def _collect_all_image_rids_from_element_tree(root: Any) -> list[str]:
    """
    a:blip 말고도 VML imagedata·일부 변형의 r:id / r:embed 등 그림 관계 ID 수집.
    (하이퍼링크 등 비이미지 r:id는 로드 후 _looks_like_raster로 걸러짐.)
    """
    out: list[str] = []
    seen: set[str] = set()
    q_embed = qn("r:embed")
    q_link = qn("r:link")
    q_rid = qn("r:id")
    direct_keys = (
        q_embed,
        q_link,
        q_rid,
        _OOXML_REL_EMBED,
        _OOXML_REL_LINK,
        _OOXML_REL_ID,
    )
    for el in root.iter():
        for k in direct_keys:
            v = el.get(k)
            if v and isinstance(v, str) and v.startswith("rId") and v not in seen:
                seen.add(v)
                out.append(v)
        loc = _xml_local_tag(el).lower()
        for ak, av in el.attrib.items():
            if not av or not isinstance(av, str) or not av.startswith("rId"):
                continue
            if av in seen:
                continue
            if ak in direct_keys or ak.endswith("}embed") or ak.endswith("}link"):
                seen.add(av)
                out.append(av)
            elif loc == "imagedata" and (ak.endswith("}id") or ak == q_rid):
                seen.add(av)
                out.append(av)
    return out


def _enrich_images_from_xml_rel_ids(
    part: Any,
    slide_part: Any,
    root_elm: Any,
    slide_w: int,
    slide_h: int,
    slide_num: int,
    used_rids: set[str],
    elements: list[dict[str, Any]],
    img_idx: int,
    image_fps: set[str] | None = None,
) -> int:
    """XML에 등장하는 모든 rId 후보로 이미지 부품 로드(도형 API 누락 보강)."""
    for r_id in _collect_all_image_rids_from_element_tree(root_elm):
        if r_id in used_rids:
            continue
        blob, ct = _load_related_image_any(part, slide_part, r_id)
        if not blob or not _looks_like_raster(blob, ct):
            continue
        used_rids.add(r_id)
        abs_l = abs_t = 0
        w_emu = max(slide_w // 5, 1)
        h_emu = max(slide_h // 5, 1)
        img_idx = _append_raster_dims(
            elements,
            blob,
            ct or "image/png",
            abs_l,
            abs_t,
            w_emu,
            h_emu,
            slide_w,
            slide_h,
            slide_num,
            img_idx,
            image_fps,
        )
    return img_idx


def _merge_slide_xml_text_into_parts(slide_elm: Any, slide_text_parts: list[str]) -> None:
    """slide XML 전체 a:t 및 (drawing이 아닌) wps/w 단락 t — 도형 순회 누락분을 plainText에 합침."""
    seen: set[str] = set()
    for part in slide_text_parts:
        for ln in part.splitlines():
            s = ln.strip()
            if s:
                seen.add(s)
    q_at = qn("a:t")
    extra: list[str] = []
    for el in slide_elm.iter():
        if el.tag != q_at or not el.text:
            continue
        t = el.text.strip()
        if t and t not in seen:
            seen.add(t)
            extra.append(t)
    for el in slide_elm.iter():
        if el.tag == q_at:
            continue
        if _xml_local_tag(el) != "t" or not el.text:
            continue
        t = el.text.strip()
        if not t or t in seen:
            continue
        seen.add(t)
        extra.append(t)
    slide_text_parts.extend(extra)


def _enrich_all_blips_in_tree(
    part: Any,
    slide_part: Any,
    root_elm: Any,
    slide_w: int,
    slide_h: int,
    slide_num: int,
    used_rids: set[str],
    elements: list[dict[str, Any]],
    img_idx: int,
    image_fps: set[str] | None = None,
) -> int:
    """
    p:pic 이외 blipFill(자동 도형·배경 등)의 a:blip까지 수집.
    관계는 주로 part에 있으나, 슬라이드 쪽 관계도 시도.
    """
    q_blip = qn("a:blip")
    q_embed = qn("r:embed")
    q_link = qn("r:link")
    for blip in root_elm.iter():
        if blip.tag != q_blip:
            continue
        for attr in (q_embed, q_link):
            r_id = blip.get(attr)
            if not r_id or r_id in used_rids:
                continue
            blob, ct = _load_related_image_any(part, slide_part, r_id)
            if not blob or not _looks_like_raster(blob, ct):
                continue
            used_rids.add(r_id)
            if _blip_under_slide_background(blip):
                abs_l, abs_t, w_emu, h_emu = 0, 0, slide_w, slide_h
            else:
                geom = _geom_from_blip_ancestor(blip)
                if geom is None:
                    abs_l, abs_t = 0, 0
                    w_emu = max(slide_w // 6, 1)
                    h_emu = max(slide_h // 6, 1)
                else:
                    abs_l, abs_t, w_emu, h_emu = geom
            img_idx = _append_raster_dims(
                elements,
                blob,
                ct or "image/png",
                abs_l,
                abs_t,
                w_emu,
                h_emu,
                slide_w,
                slide_h,
                slide_num,
                img_idx,
                image_fps,
            )
    return img_idx


def _enrich_slide_orphan_image_rels(
    slide_part: Any,
    slide_w: int,
    slide_h: int,
    slide_num: int,
    used_rids: set[str],
    elements: list[dict[str, Any]],
    img_idx: int,
    image_fps: set[str] | None = None,
) -> int:
    """슬라이드 .rels에만 있는 이미지(비표준 XML 등) 보강. 레이아웃/마스터는 호출하지 않음."""
    try:
        rels = getattr(slide_part, "rels", None)
        if rels is None:
            return img_idx
        for rel in rels.values():
            rt = getattr(rel, "reltype", "") or ""
            if "relationships/image" not in rt:
                continue
            rid = getattr(rel, "rId", None)
            if not rid or rid in used_rids:
                continue
            try:
                tp = rel.target_part
                blob = getattr(tp, "blob", None)
                ct = (getattr(tp, "content_type", None) or "").strip() or None
            except Exception:
                continue
            if not blob or not _looks_like_raster(blob, ct):
                continue
            used_rids.add(rid)
            img_idx = _append_raster_dims(
                elements,
                blob,
                ct or "image/png",
                0,
                0,
                max(slide_w // 4, 1),
                max(slide_h // 4, 1),
                slide_w,
                slide_h,
                slide_num,
                img_idx,
                image_fps,
            )
    except Exception:
        pass
    return img_idx


def _enrich_oxml_pics(
    part: Any,
    slide_part: Any,
    root_elm: Any,
    slide_w: int,
    slide_h: int,
    slide_num: int,
    used_rids: set[str],
    elements: list[dict[str, Any]],
    img_idx: int,
    image_fps: set[str] | None = None,
) -> int:
    """
    python-pptx shape 트리 밖(mc:AlternateContent 등)의 p:pic 보강.
    part.related_part으로 해당 rId 해석.
    """
    for pic in root_elm.iter():
        if pic.tag != qn("p:pic"):
            continue
        for r_id in _pic_blip_rel_ids(pic):
            if r_id in used_rids:
                continue
            blob, ct = _load_related_image_any(part, slide_part, r_id)
            if not blob or not _looks_like_raster(blob, ct):
                continue
            used_rids.add(r_id)
            geom = _pic_geom_emu(pic)
            if geom is None:
                abs_l = abs_t = 0
                w_emu = max(slide_w // 20, 1)
                h_emu = max(slide_h // 20, 1)
            else:
                abs_l, abs_t, w_emu, h_emu = geom
            img_idx = _append_raster_dims(
                elements,
                blob,
                ct or "image/png",
                abs_l,
                abs_t,
                w_emu,
                h_emu,
                slide_w,
                slide_h,
                slide_num,
                img_idx,
                image_fps,
            )
    return img_idx


def _content_type_from_zip_path(name: str) -> str:
    n = name.lower()
    if n.endswith(".png"):
        return "image/png"
    if n.endswith(".jpg") or n.endswith(".jpeg"):
        return "image/jpeg"
    if n.endswith(".gif"):
        return "image/gif"
    if n.endswith(".webp"):
        return "image/webp"
    if n.endswith(".bmp"):
        return "image/bmp"
    if n.endswith(".tif") or n.endswith(".tiff"):
        return "image/tiff"
    if n.endswith(".emf"):
        return "image/x-emf"
    if n.endswith(".wmf"):
        return "image/x-wmf"
    if n.endswith(".svg"):
        return "image/svg+xml"
    return "application/octet-stream"


def _zip_resolve_part_target(part_dir: str, target: str) -> str:
    if target.startswith("/"):
        return target.lstrip("/").replace("\\", "/")
    return posixpath.normpath(posixpath.join(part_dir, target)).replace("\\", "/")


def _zip_iter_image_targets_from_rels(
    zf: zipfile.ZipFile,
    rels_internal_path: str,
    part_dir_for_targets: str,
) -> list[str]:
    """rels 파일에서 내부 이미지 타깃 경로(zip 내) 목록."""
    if rels_internal_path not in zf.namelist():
        return []
    out: list[str] = []
    try:
        rels_root = ET.fromstring(zf.read(rels_internal_path))
        for el in rels_root:
            loc = el.tag.split("}")[-1] if "}" in el.tag else el.tag
            if loc != "Relationship":
                continue
            typ = (el.get("Type") or "").lower()
            if "relationships/image" not in typ and "/image" not in typ:
                continue
            if el.get("TargetMode") == "External":
                continue
            target = el.get("Target")
            if not target:
                continue
            zip_path = _zip_resolve_part_target(part_dir_for_targets, target)
            if zip_path in zf.namelist():
                out.append(zip_path)
    except Exception:
        pass
    return out


def _slide_zip_package_enrich(
    zf: zipfile.ZipFile,
    slide_num: int,
    slide_w: int,
    slide_h: int,
    elements: list[dict[str, Any]],
    slide_text_parts: list[str],
    img_idx: int,
    image_fps: set[str] | None = None,
) -> tuple[int, dict[str, Any]]:
    """
    OOXML ZIP에서 slideN 관계·레이아웃·마스터 rels의 image Target까지 따라가 media 로드.
    """
    stats: dict[str, Any] = {
        "slideXmlInZip": False,
        "slideRelsInZip": False,
        "zipImageRelCount": 0,
        "zipImagesEmitted": 0,
        "zipLayoutMasterPartsFollowed": 0,
        "zipTextFragmentsAdded": 0,
    }
    slide_xml = f"ppt/slides/slide{slide_num}.xml"
    rels_xml = f"ppt/slides/_rels/slide{slide_num}.xml.rels"
    base_dir = "ppt/slides"
    emitted_paths: set[str] = set()

    def try_emit_zip_paths(paths: list[str]) -> None:
        nonlocal img_idx
        for zip_path in paths:
            if zip_path in emitted_paths or zip_path not in zf.namelist():
                continue
            try:
                blob = zf.read(zip_path)
            except Exception:
                continue
            ct = _content_type_from_zip_path(zip_path)
            if not _looks_like_raster(blob, ct):
                continue
            emitted_paths.add(zip_path)
            prev_idx = img_idx
            img_idx = _append_raster_dims(
                elements,
                blob,
                ct,
                0,
                0,
                max(slide_w // 4, 1),
                max(slide_h // 4, 1),
                slide_w,
                slide_h,
                slide_num,
                img_idx,
                image_fps,
            )
            if img_idx > prev_idx:
                stats["zipImagesEmitted"] += 1

    if slide_xml in zf.namelist():
        stats["slideXmlInZip"] = True
    if rels_xml in zf.namelist():
        stats["slideRelsInZip"] = True
        try:
            rels_root = ET.fromstring(zf.read(rels_xml))
            slide_image_targets: list[str] = []
            layout_master_xmls: list[str] = []
            for el in rels_root:
                loc = el.tag.split("}")[-1] if "}" in el.tag else el.tag
                if loc != "Relationship":
                    continue
                typ = (el.get("Type") or "").lower()
                target = el.get("Target")
                if not target or el.get("TargetMode") == "External":
                    continue
                if "relationships/image" in typ or "/image" in typ:
                    stats["zipImageRelCount"] += 1
                    zp = _zip_resolve_part_target(base_dir, target)
                    if zp in zf.namelist():
                        slide_image_targets.append(zp)
                elif "slidelayout" in typ or "slidemaster" in typ:
                    lpath = _zip_resolve_part_target(base_dir, target)
                    if lpath in zf.namelist() and lpath.endswith(".xml"):
                        layout_master_xmls.append(lpath)

            try_emit_zip_paths(slide_image_targets)

            seen_lm: set[str] = set()
            for lm_xml in layout_master_xmls:
                if lm_xml in seen_lm:
                    continue
                seen_lm.add(lm_xml)
                stats["zipLayoutMasterPartsFollowed"] += 1
                lm_dir = posixpath.dirname(lm_xml)
                lm_base = posixpath.basename(lm_xml)
                lm_rels = posixpath.join(lm_dir, "_rels", lm_base + ".rels")
                nested = _zip_iter_image_targets_from_rels(zf, lm_rels, lm_dir)
                stats["zipImageRelCount"] += len(nested)
                try_emit_zip_paths(nested)
        except Exception:
            pass

    if slide_xml in zf.namelist():
        try:
            sroot = ET.fromstring(zf.read(slide_xml))
            seen: set[str] = set()
            for part in slide_text_parts:
                for ln in part.splitlines():
                    s = ln.strip()
                    if s:
                        seen.add(s)
            drawing_t = "{http://schemas.openxmlformats.org/drawingml/2006/main}t"
            extra: list[str] = []
            for el in sroot.iter():
                if el.text is None:
                    continue
                t = el.text.strip()
                if not t or t in seen:
                    continue
                tag = el.tag
                if tag == drawing_t:
                    seen.add(t)
                    extra.append(t)
                    continue
                local = tag.split("}")[-1] if "}" in tag else tag
                if local == "t" and "wordprocessing" in tag.lower():
                    seen.add(t)
                    extra.append(t)
            stats["zipTextFragmentsAdded"] = len(extra)
            slide_text_parts.extend(extra)
        except Exception:
            pass

    return img_idx, stats


def _load_theme_elm_from_zip(content: bytes) -> Any:
    """
    PPTX ZIP에서 ppt/theme/theme*.xml 을 직접 읽어 ElementTree 엘리먼트를 반환.
    python-pptx part 체인과 무관하게 동작하므로 SmartArt/그룹 도형에서도 안정적.
    """
    try:
        with zipfile.ZipFile(io.BytesIO(content)) as z:
            theme_names = sorted(
                n for n in z.namelist()
                if n.startswith("ppt/theme/") and n.endswith(".xml")
            )
            if not theme_names:
                return None
            return ET.fromstring(z.read(theme_names[0]))
    except Exception:
        return None


def parse_presentation(content: bytes) -> dict[str, Any]:
    prs = Presentation(io.BytesIO(content))
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    # 테마 XML을 ZIP에서 직접 추출 (part 체인보다 안정적)
    global_theme_elm = _load_theme_elm_from_zip(content)

    all_plain: list[str] = []
    slides_out: list[dict[str, Any]] = []

    zf: zipfile.ZipFile | None = None
    try:
        zf = zipfile.ZipFile(io.BytesIO(content))
    except Exception:
        zf = None

    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        elements: list[dict[str, Any]] = []
        slide_text_parts: list[str] = []
        img_idx = 0
        image_fps: set[str] = set()

        slide_part = slide.part
        used_blip_rids: set[str] = set()

        # _iter_shapes_placed가 이제 슬라이드 절대 EMU를 직접 반환
        for shape, abs_l, abs_t in _iter_shapes_paint_order(slide):

            st = _shape_type_safe(shape)
            skip_blip_after_picture = False
            if st in (
                MSO_SHAPE_TYPE.PICTURE,
                MSO_SHAPE_TYPE.LINKED_PICTURE,
            ):
                blob: bytes | None = None
                ct: str | None = None
                try:
                    blob = shape.image.blob  # type: ignore[attr-defined]
                    ct = shape.image.content_type  # type: ignore[attr-defined]
                except Exception:
                    blob, ct = None, None
                if blob and _looks_like_raster(blob, ct):
                    for rid in _blip_rel_ids(shape.element):
                        used_blip_rids.add(rid)
                    img_idx = _append_raster_element(
                        elements,
                        blob,
                        ct or "image/png",
                        abs_l,
                        abs_t,
                        shape,
                        slide_w,
                        slide_h,
                        slide_num,
                        img_idx,
                        image_fps,
                    )
                    skip_blip_after_picture = True

            if not skip_blip_after_picture:
                seen_rid: set[str] = set()
                for r_id in _blip_rel_ids(shape.element):
                    if r_id in seen_rid:
                        continue
                    seen_rid.add(r_id)
                    blob, ct = _load_related_image_any(shape.part, slide_part, r_id)
                    if blob and _looks_like_raster(blob, ct):
                        used_blip_rids.add(r_id)
                        img_idx = _append_raster_element(
                            elements,
                            blob,
                            ct or "image/png",
                            abs_l,
                            abs_t,
                            shape,
                            slide_w,
                            slide_h,
                            slide_num,
                            img_idx,
                            image_fps,
                        )

            # 표 도형 → type: "table" 구조 추출
            if getattr(shape, "has_table", False):
                try:
                    rows_data: list[list[str]] = []
                    for row in shape.table.rows:
                        cells_data: list[str] = []
                        for cell in row.cells:
                            try:
                                txt_cell = (cell.text_frame.text or "").strip()
                            except Exception:
                                txt_cell = (getattr(cell, "text", "") or "").strip()
                            cells_data.append(txt_cell)
                        rows_data.append(cells_data)
                    if rows_data:
                        flat_text = "\n".join(" | ".join(r) for r in rows_data)
                        slide_text_parts.append(flat_text)
                        elements.append({
                            "type": "table",
                            "rows": rows_data,
                            "style": {
                                "left": _pct(abs_l, slide_w),
                                "top": _pct(abs_t, slide_h),
                                "width": _pct(shape.width, slide_w),
                                "height": _pct(shape.height, slide_h),
                            },
                        })
                        continue  # 표는 text 경로로 다시 처리하지 않음
                except Exception:
                    pass

            # 채워진 도형·커넥터 추출 (텍스트 유무와 무관하게 배경 fill 시도)
            txt = _shape_text(shape)
            if not txt.strip():
                _try_extract_connector_shape(
                    shape, abs_l, abs_t, slide_w, slide_h, elements, global_theme_elm
                )
            # 텍스트가 있어도 배경 fill이 있으면 shape 요소로 추가 (조직도 노드 등)
            _try_extract_filled_shape(
                shape, abs_l, abs_t, slide_w, slide_h, elements, global_theme_elm
            )
            if txt.strip():
                slide_text_parts.append(txt.strip())
            if txt.strip():
                left_pct = _pct(abs_l, slide_w)
                top_pct = _pct(abs_t, slide_h)
                width_pct = _pct(shape.width, slide_w)
                height_pct = _pct(shape.height, slide_h)
                text_style = _extract_text_style(shape, slide_w)
                para_styles = _extract_paragraph_styles(shape, slide_w)
                style: dict[str, Any] = {
                    "left": left_pct,
                    "top": top_pct,
                    "width": width_pct,
                    "height": height_pct,
                }
                style.update(text_style)
                elem: dict[str, Any] = {
                    "type": "text",
                    "content": txt,
                    "style": style,
                }
                # 단락별 스타일이 실질적으로 다양할 때만 포함 (단일 스타일이면 불필요)
                non_empty = [p for p in para_styles if p]
                unique_styles = {tuple(sorted(p.items())) for p in non_empty}
                if len(unique_styles) > 1:
                    elem["paragraphStyles"] = para_styles
                elements.append(elem)

        try:
            lo = slide.slide_layout
            ma = lo.slide_master
            oxml_roots: tuple[tuple[Any, Any], ...] = (
                (slide_part, slide._element),
                (lo.part, lo._element),
                (ma.part, ma._element),
            )
        except Exception:
            oxml_roots = ((slide_part, slide._element),)
        for oxml_part, root_elm in oxml_roots:
            img_idx = _enrich_oxml_pics(
                oxml_part,
                slide_part,
                root_elm,
                slide_w,
                slide_h,
                slide_num,
                used_blip_rids,
                elements,
                img_idx,
                image_fps,
            )
        for oxml_part, root_elm in oxml_roots:
            img_idx = _enrich_all_blips_in_tree(
                oxml_part,
                slide_part,
                root_elm,
                slide_w,
                slide_h,
                slide_num,
                used_blip_rids,
                elements,
                img_idx,
                image_fps,
            )
        img_idx = _enrich_images_from_xml_rel_ids(
            slide_part,
            slide_part,
            slide._element,
            slide_w,
            slide_h,
            slide_num,
            used_blip_rids,
            elements,
            img_idx,
            image_fps,
        )
        img_idx = _enrich_slide_orphan_image_rels(
            slide_part,
            slide_w,
            slide_h,
            slide_num,
            used_blip_rids,
            elements,
            img_idx,
            image_fps,
        )

        _enrich_tables_from_slide_oxml(
            slide._element,
            slide_w,
            slide_h,
            elements,
            slide_text_parts,
        )

        _merge_slide_xml_text_into_parts(slide._element, slide_text_parts)

        extract_stats: dict[str, Any] = {
            "zipPackageOk": zf is not None,
            "parserApiBuild": PARSER_API_BUILD,
        }
        if zf is not None:
            img_idx, zst = _slide_zip_package_enrich(
                zf,
                slide_num,
                slide_w,
                slide_h,
                elements,
                slide_text_parts,
                img_idx,
                image_fps,
            )
            extract_stats.update(zst)

        plain_slide = "\n".join(slide_text_parts)
        if plain_slide:
            all_plain.append(plain_slide)
        slides_out.append(
            {
                "slideNumber": slide_num,
                "elements": elements,
                "plainText": plain_slide,
                "extractStats": extract_stats,
            }
        )

    if zf is not None:
        try:
            zf.close()
        except Exception:
            pass

    full_plain = "\n\n".join(all_plain)
    words = re.findall(r"[\w가-힣]+", full_plain)
    freq: dict[str, int] = {}
    for w in words:
        if len(w) < 2:
            continue
        freq[w.lower()] = freq.get(w.lower(), 0) + 1
    top_tags = sorted(freq, key=lambda k: freq[k], reverse=True)[:5]

    title_guess = ""
    if slides_out and slides_out[0].get("elements"):
        for el in slides_out[0]["elements"]:
            if el.get("type") == "text" and el.get("content"):
                title_guess = (el["content"] or "").split("\n")[0].strip()[:120]
                break

    return {
        "slides": slides_out,
        "plainText": full_plain,
        "meta": {
            "title": title_guess or "Untitled lecture",
            "description": "",
            "tags": top_tags,
            "parserApiBuild": PARSER_API_BUILD,
            "slideSizeEmu": {
                "width": int(prs.slide_width),
                "height": int(prs.slide_height),
            },
        },
    }


@app.post("/api/parse-pptx")
async def parse_pptx(file: UploadFile = File(...)):
    if not file.filename or not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Expected .pptx file")
    content = await file.read()
    if len(content) == 0:
        raise HTTPException(status_code=400, detail="Empty file")
    try:
        payload = parse_presentation(content)
        slides = payload.get("slides") or []

        import asyncio, logging
        logger = logging.getLogger("pptx_raster")

        # PowerPoint COM은 STA COM이라 async 이벤트 루프 스레드에서 직접 호출하면
        # 아파트먼트 충돌이 발생함 → 전용 스레드에서 실행
        raster_urls, raster_meta = await asyncio.to_thread(
            render_slide_rasters_ppt, content, len(slides)
        )
        ppt_status = raster_meta.get("status")
        ppt_reason = raster_meta.get("reason", "")
        logger.warning("[raster] ppt-com status=%s reason=%s", ppt_status, ppt_reason)

        if ppt_status not in ("ok", "partial"):
            # PPT COM 실패 이유를 meta에 보존하고 LO로 재시도
            lo_urls, lo_meta = await asyncio.to_thread(
                render_slide_rasters_jpeg, content, len(slides)
            )
            lo_meta["pptComFallbackReason"] = ppt_reason or "ppt-com status=" + str(ppt_status)
            lo_meta["engine"] = "libreoffice"
            raster_urls, raster_meta = lo_urls, lo_meta
        
        for slide_dict, url in zip(slides, raster_urls):
            if url:
                slide_dict["rasterPreview"] = url
        meta = payload.setdefault("meta", {})
        meta["slideRaster"] = raster_meta
        return JSONResponse(
            content=payload,
            headers={"X-PPTX-Parser-Build": PARSER_API_BUILD},
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e)) from e


@app.get("/health")
def health():
    # 브라우저가 JSON을 오래 캐시하면 parserApiBuild 갱신이 안 보일 수 있음
    return JSONResponse(
        content={"status": "ok", "parserApiBuild": PARSER_API_BUILD},
        headers={"Cache-Control": "no-store, max-age=0"},
    )


@app.get("/pptx-parser/build")
def parser_build_info():
    """같은 호스트에 다른 앱이 있을 때 /health만으로 헷갈리면 이 경로로 판별."""
    return JSONResponse(
        content={
            "parserApiBuild": PARSER_API_BUILD,
            "mainFile": __file__,
        },
        headers={"Cache-Control": "no-store, max-age=0"},
    )
