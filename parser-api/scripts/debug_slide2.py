"""One-off: inspect slide 2 shapes in Desktop pptx files."""
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from app.main import _blip_rel_ids, _iter_shapes_placed, _shape_text


def main() -> None:
    import os

    desktop = Path(r"c:/Users/USER/Desktop")
    names = [n for n in os.listdir(desktop) if n.lower().endswith(".pptx")]
    paths = sorted(
        (desktop / n for n in names),
        key=lambda p: p.stat().st_size,
        reverse=True,
    )[:8]
    for p in paths:
        print("===", p.name, p.stat().st_size)
        try:
            prs = Presentation(str(p))
        except Exception as e:
            print("  open ERR", e)
            continue
        if len(prs.slides) < 2:
            print("  <2 slides")
            continue
        slide = prs.slides[1]
        for idx, (sh, _gl, _gt) in enumerate(_iter_shapes_placed(slide.shapes)):
            et = sh.element.tag
            local = et.split("}")[-1] if "}" in et else et
            st = sh.shape_type
            blips = _blip_rel_ids(sh.element)
            for node in sh.element.iter():
                if node.tag == qn("a:blip"):
                    print(
                        "    blip attrs",
                        {k: v for k, v in node.attrib.items()},
                    )
                    break
            ht = getattr(sh, "has_table", False)
            raw = _shape_text(sh) or ""
            txt_preview = raw[:100].replace("\n", " ")
            pic = st in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE)
            print(
                f"  {idx}: st={st} xml={local} has_table={ht} "
                f"blip_embeds={blips} picture_enum={pic} txt={txt_preview!r}"
            )


if __name__ == "__main__":
    main()
